import fitz


def read_pdf(uploaded_file):
    """Open the uploaded PDF once."""

    pdf_bytes = uploaded_file.getvalue()

    pdf_document = fitz.open(
        stream=pdf_bytes,
        filetype="pdf"
    )

    return pdf_document


def extract_text_from_pdf(uploaded_file):
    """
    Extract text from an uploaded PDF.

    BUGFIX: the previous version used page.get_text() (default
    "text" mode) and simply concatenated every page's text with
    "+=". Two problems with that:

    1. get_text()'s default reading order is based on the PDF's
       internal draw order, not visual position - for PDFs
       exported from PowerPoint (very common for lecture slides),
       small text boxes like a footer ("Presented by: Dr. ...
       Assistant Professor  Department of ...") are often drawn
       as several separate, closely-spaced runs. Those runs can
       get concatenated with NO whitespace between them ("Bhopale.
       Assistant ProfessorDepartment of..."), which then defeats
       both heading detection (it's not a clean short line anymore)
       and boilerplate-repeat detection (no two pages produce an
       identical line to match against).
    2. Concatenating page texts with a bare "+=" means the last
       line of one page can run directly into the first line of
       the next page if the page's own text doesn't end in a
       newline.

    Fix: extract text BLOCK BY BLOCK (each visually distinct text
    box becomes its own block), sort blocks top-to-bottom /
    left-to-right, and join them with explicit newlines - so every
    visually separate piece of text (including a footer credit
    line) always ends up on its own clean line. Pages are joined
    with a form-feed character ("\\x0c"), which Python's
    str.splitlines() already treats as a line boundary, so nothing
    downstream needs to change to benefit from the page boundary -
    it also lets boilerplate-detection code reason about "how many
    pages" a repeated line appeared on if it ever needs to.
    """

    pdf_document = read_pdf(uploaded_file)

    page_texts = []

    for page in pdf_document:

        blocks = page.get_text("blocks")

        # Sort top-to-bottom, then left-to-right, so reading
        # order matches how a human would actually read the page.
        blocks = sorted(
            blocks,
            key=lambda block: (round(block[1], 1), block[0])
        )

        block_lines = [
            block[4].strip()
            for block in blocks
            if block[4] and block[4].strip()
        ]

        page_texts.append(
            "\n".join(block_lines)
        )

    pdf_document.close()

    return "\x0c".join(page_texts)


def get_pdf_page_count(uploaded_file):
    """Return the number of pages in a PDF."""

    pdf_document = read_pdf(uploaded_file)

    page_count = len(pdf_document)

    pdf_document.close()

    return page_count


# ============================================================
# MULTI-FORMAT INGESTION (DOCX / PPTX / TXT)
#
# The PS asks for books, textbooks, PDFs, notes, DOC/DOCX,
# PPT/PPTX and other text-based material. These additions sit
# next to the existing PDF functions above (which are untouched)
# and give app.py one dispatcher to call regardless of file type.
# ============================================================

def extract_text_from_docx(uploaded_file):
    """
    Extract text from an uploaded .docx file, including paragraph
    text and any text inside tables (tables are common in lecture
    notes and are otherwise silently skipped by a plain paragraph
    loop).
    """

    import docx

    document = docx.Document(uploaded_file)

    parts = [
        paragraph.text
        for paragraph in document.paragraphs
        if paragraph.text and paragraph.text.strip()
    ]

    for table in document.tables:
        for row in table.rows:
            row_text = " | ".join(
                cell.text.strip()
                for cell in row.cells
                if cell.text and cell.text.strip()
            )
            if row_text:
                parts.append(row_text)

    return "\n".join(parts)


def extract_text_from_pptx(uploaded_file):
    """
    Extract text from an uploaded .pptx file: every text-bearing
    shape on every slide, plus speaker notes (often where the real
    explanation lives on a slide deck), one slide per block so the
    downstream heading/section detection still sees natural breaks.
    """

    from pptx import Presentation

    presentation = Presentation(uploaded_file)

    slide_blocks = []

    for slide_number, slide in enumerate(presentation.slides, start=1):

        lines = [f"Slide {slide_number}"]

        for shape in slide.shapes:

            if shape.has_text_frame and shape.text_frame.text.strip():
                lines.append(shape.text_frame.text.strip())

            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(
                        cell.text.strip()
                        for cell in row.cells
                        if cell.text and cell.text.strip()
                    )
                    if row_text:
                        lines.append(row_text)

        if slide.has_notes_slide:
            notes_text = (
                slide.notes_slide.notes_text_frame.text
                if slide.notes_slide.notes_text_frame
                else ""
            )
            if notes_text and notes_text.strip():
                lines.append(f"Notes: {notes_text.strip()}")

        if len(lines) > 1:
            slide_blocks.append("\n".join(lines))

    return "\n\n".join(slide_blocks)


def extract_text_from_txt(uploaded_file):
    """Extract text from an uploaded .txt or .md file."""

    raw_bytes = uploaded_file.getvalue()

    try:
        return raw_bytes.decode("utf-8")
    except UnicodeDecodeError:
        return raw_bytes.decode("utf-8", errors="ignore")


def get_file_extension(uploaded_file):
    """Return the lowercase extension of an uploaded file, no dot."""

    name = getattr(uploaded_file, "name", "") or ""

    if "." not in name:
        return ""

    return name.rsplit(".", 1)[-1].lower()


def extract_text_from_upload(uploaded_file):
    """
    Single entry point for app.py: detect the uploaded file's type
    and route it to the right extractor.

    Returns (text, page_count). page_count is a real page count
    for PDFs (existing behavior, unchanged) and "N/A" for formats
    that don't have a natural page concept.
    """

    extension = get_file_extension(uploaded_file)

    if extension == "pdf":
        return extract_text_from_pdf(uploaded_file), get_pdf_page_count(uploaded_file)

    if extension in ("docx", "doc"):
        return extract_text_from_docx(uploaded_file), "N/A"

    if extension in ("pptx", "ppt"):
        return extract_text_from_pptx(uploaded_file), "N/A"

    if extension in ("txt", "md"):
        return extract_text_from_txt(uploaded_file), "N/A"

    raise ValueError(
        f"Unsupported file type: .{extension or 'unknown'}. "
        "Please upload a PDF, DOCX, PPTX, or TXT file."
    )